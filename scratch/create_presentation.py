import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(248, 246, 240)       # #F8F6F0 Vintage Cream Paper
    COLOR_INK = RGBColor(26, 26, 26)         # #1A1A1A Deep Ink
    COLOR_YELLOW = RGBColor(255, 222, 89)    # #FFDE59 Acid Yellow
    COLOR_PEACH = RGBColor(255, 235, 225)    # #FFEBE1 Soft Peach
    COLOR_CARD = RGBColor(255, 255, 255)     # #FFFFFF White
    COLOR_MUTED = RGBColor(100, 100, 100)    # Gray

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="FRAMES APP ARCHITECTURE"):
        # Header category text
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_MUTED

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_INK

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    # Decorative Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_INK
    card1.line.width = Pt(2)

    # Title inside card
    tbox = slide1.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.333), Inches(3.0))
    tf = tbox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "📸 FRAMES MOBILE APP"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_MUTED
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "Tactile Polaroid Social Scrapbook & Print-on-Demand Architecture"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_INK
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "Complete Project Study Guide: System Architecture, Supabase Database, AdMob & POD API Triggers"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_INK
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 2: Executive Overview
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "1. Executive Product Vision & Key Differentiators")

    features = [
        ("📸 Tactile Polaroid Scrapbook", "Users capture daily moments formatted as realistic physical polaroids with film grain, brand avatars, tape pins, and interactive reactions."),
        ("📦 $4.99 Daily Memory Packs", "Automated Print-on-Demand (POD) integration that compiles a user's daily digital frames into physical glossy prints delivered to their doorstep."),
        ("📢 AdMob Placement Containers", "Minimalist ad containers formatted like sponsored polaroids in the feed, supporting Native Banner & Rewarded Video ads."),
        ("🔔 Real-Time Mobile Push Engine", "OS-level heads-up notification banners for friend requests, direct chat messages, and live order tracking updates.")
    ]

    for i, (feat_title, feat_desc) in enumerate(features):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.8)
        top = Inches(1.7 + row * 2.6)

        c = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.5), Inches(2.3))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_INK
        c.line.width = Pt(1.5)

        tb = slide2.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.1), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = feat_title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_INK

        p2 = tf.add_paragraph()
        p2.text = feat_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture & Tech Stack
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "2. Full Technical Architecture & Component Stack")

    stack_items = [
        ("Mobile App", "React Native + Expo Router v4", "Cross-platform iOS/Android app & Web export"),
        ("State Management", "Zustand (appStore.ts)", "Ultra-fast global client state management"),
        ("Database & Auth", "Supabase (PostgreSQL + RLS)", "User auth, daily frames DB, Row-Level Security"),
        ("Monetization", "Google AdMob Console", "Banner & Rewarded Video ad placements"),
        ("Push Engine", "Expo Push Service", "Android high-importance OS notification channels"),
        ("Print Lab API", "Prodigi / Printo POD API", "Automated order creation & shipping webhook dispatches")
    ]

    for i, (layer, tech, purpose) in enumerate(stack_items):
        top = Inches(1.6 + i * 0.9)

        c = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(0.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_INK
        c.line.width = Pt(1)

        # Layer Name
        tb1 = slide3.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(2.8), Inches(0.5))
        p1 = tb1.text_frame.paragraphs[0]
        p1.text = layer
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_INK

        # Tech Name
        tb2 = slide3.shapes.add_textbox(Inches(4.0), top + Inches(0.15), Inches(3.8), Inches(0.5))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = tech
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(180, 100, 0)

        # Purpose
        tb3 = slide3.shapes.add_textbox(Inches(8.0), top + Inches(0.15), Inches(4.3), Inches(0.5))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = purpose
        p3.font.size = Pt(13)
        p3.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 4: Backend Database Schema
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "3. Backend PostgreSQL Database Schema & Security")

    tables = [
        ("profiles Table", "Stores user identities, display names, avatars, and push tokens.", "id (UUID PK), username, display_name, avatar_url, push_token"),
        ("frames Table", "Stores daily captured polaroid photo URLs and captions.", "id (UUID PK), user_id (FK), image_url, caption, created_at"),
        ("print_orders Table", "Stores physical $4.99 Daily Memory Pack print orders.", "id (UUID PK), user_id, photos (JSONB), shipping_address, tracking_number")
    ]

    for i, (tbl_title, tbl_desc, tbl_schema) in enumerate(tables):
        top = Inches(1.6 + i * 1.8)

        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.6))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_INK
        c.line.width = Pt(1.5)

        tb = slide4.shapes.add_textbox(Inches(1.1), top + Inches(0.15), Inches(11.133), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = tbl_title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_INK

        p2 = tf.add_paragraph()
        p2.text = tbl_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = f"Columns: {tbl_schema}"
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(30, 120, 30)
        p3.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 5: API Triggers & Print-on-Demand Flow
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "4. $4.99 Daily Print Pack API Trigger Pipeline")

    steps = [
        ("Step 1: Checkout Modal", "User opens OrderDailyPackModal.tsx from Archive/Timeline and submits address."),
        ("Step 2: DB Order Creation", "createRemotePrintOrder() saves photo array & shipping address into Supabase."),
        ("Step 3: POD Lab Dispatch", "sendOrderToPrintPartner() compiles photo URLs & dispatches print job to Prodigi API."),
        ("Step 4: Push Notification", "triggerLocalPushNotification() sends instant OS push banner with tracking ID.")
    ]

    for i, (step_title, step_desc) in enumerate(steps):
        left = Inches(0.8 + i * 2.95)
        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(2.8), Inches(4.2))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_INK
        c.line.width = Pt(1.5)

        tb = slide5.shapes.add_textbox(left + Inches(0.15), Inches(2.2), Inches(2.5), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = step_title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_INK

        p2 = tf.add_paragraph()
        p2.text = step_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 6: Code Structure & How to Run
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "5. Project Execution & Deployment Commands")

    cmds = [
        ("Local Web App", "npx expo start --web", "Runs local web app server on port 8081"),
        ("Expo Go (Mobile)", "npx expo start --port 8082", "Connects to Expo Go mobile app"),
        ("Vercel Web Build", "npx pnpm --filter @frames/mobile build", "Compiles SPA export into dist/ folder"),
        ("EAS OTA Update", "npx eas-cli update --branch preview", "Publishes instant Over-The-Air JS updates")
    ]

    for i, (cmd_title, cmd_code, cmd_desc) in enumerate(cmds):
        top = Inches(1.6 + i * 1.35)

        c = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.15))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_INK
        c.line.width = Pt(1.5)

        tb = slide6.shapes.add_textbox(Inches(1.1), top + Inches(0.12), Inches(11.133), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"{cmd_title}:  {cmd_code}"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_INK

        p2 = tf.add_paragraph()
        p2.text = cmd_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(4)

    output_path = r"c:\Users\karth\OneDrive\Desktop\AI generated frames\Frames_App_Architecture_Presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    create_deck()
